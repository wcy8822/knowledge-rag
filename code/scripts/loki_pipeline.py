#!/usr/bin/env python3
"""
Loki — 本地知识向量化流水线 v2
覆盖: Obsidian MD / PDF / Office(docx/xlsx/pptx) / MySQL DDL
特性: 断点续跑 · 自动跳过已处理 · 单文件错误不中断 · 结构化日志
"""

import os, sys, time, json, hashlib, logging, traceback, fcntl, gc
from datetime import datetime
from pathlib import Path

import pymysql
import chromadb
import torch
from sentence_transformers import SentenceTransformer

from loki_state import StateV2, build_chroma_loader
from loki_ddl_filter import is_excluded_table
from loki_text_extract import extract_for_embed, is_code_path
from loki_scan_filter import is_excluded_path, is_low_value_file

# ========== 路径配置 ==========
BGE_PATH     = "/Users/didi/Work/projects/knowledge-rag-知识库/bge-m3-model/bge-m3/BAAI/bge-m3"
CHROMA_DIR   = "/Users/didi/Work/data/vectors/data/chroma-doc-knowledge-bge"
LOG_DIR      = Path("/Users/didi/Work/projects/knowledge-rag-知识库/code/scripts/logs")
STATE_FILE   = LOG_DIR / "loki_state.json"  # 已处理文件的fingerprint

SCAN_DIRS = [
    "/Users/didi/Work/docs/obsidian-vault/obsidian",   # Obsidian笔记
    "/Users/didi/Work/docs",                            # 其他文档
    "/Users/didi/Work/projects",                        # 项目文档
    "/Users/didi/.allin_ai_safe/docs",                  # 知识库算法文档
]
EXCLUDE_DIRS = {'.obsidian', 'Templates', 'node_modules', '__pycache__',
                '.git', 'venv', 'env', 'bge-m3-model', 'archives', 'chroma',
                '人事-简历', '.pytest_cache', 'logs'}
SUPPORTED_EXT = {'.md', '.pdf', '.docx', '.xlsx', '.pptx', '.txt', '.sql',
                  '.py', '.sh', '.yaml', '.yml'}

COLLECTION_DOCS   = "doc_knowledge_bge_m3"
COLLECTION_DDL    = "ddl_schema_bge_m3"
COLLECTION_CHUNKS = "doc_knowledge_chunks"
LOCK_FILE = Path("/Users/didi/Work/data/vectors/data/chroma-doc-knowledge-bge/.loki_write.lock")

def _get_mysql_password():
    val = os.environ.get('MYSQL_ROOT_PASSWORD')
    if val:
        return val
    secrets_file = Path.home() / ".secrets.env"
    if secrets_file.exists():
        for line in secrets_file.read_text().splitlines():
            if line.startswith("export MYSQL_ROOT_PASSWORD="):
                return line.split("=", 1)[1].strip().strip('"').strip("'")
    return ""

DB_CONFIG = dict(host='localhost', user='root',
                 password=_get_mysql_password(),
                 charset='utf8mb4', connect_timeout=5)

BATCH_SIZE    = 12             # M5+24G：4→12，摊平 GPU batch overhead
MAX_CHARS     = 2000
SLEEP_BATCH   = 0.2            # MPS leak 已修，0.2s 留一点 GC 余量即可
_MODEL_RELOAD = False           # MPS context 泄漏已修复，不再需要周期性重建 model


def _free_mps_memory(verbose: bool = False) -> None:
    """每批后必调：释放 PyTorch MPS 缓存 + Python GC。
    防 PyTorch issue #154329 长跑泄漏。"""
    gc.collect()
    if torch.backends.mps.is_available():
        try:
            torch.mps.empty_cache()
        except Exception as e:
            if verbose:
                log.warning(f"  torch.mps.empty_cache 失败: {e}")



# ========== 日志 ==========
LOG_DIR.mkdir(parents=True, exist_ok=True)
log_file = LOG_DIR / f"loki_{datetime.now().strftime('%Y%m%d_%H%M%S')}.log"
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s %(levelname)s %(message)s',
    handlers=[
        logging.FileHandler(log_file, encoding='utf-8'),
        logging.StreamHandler(sys.stdout),
    ]
)
log = logging.getLogger('loki')

# ========== 状态持久化 (v2: dict[path → fp]，详见 loki_state.py) ==========
def load_state() -> StateV2:
    """加载 v2 state；若读到 v1 List 自动从 ChromaDB 迁移。"""
    return StateV2.load(STATE_FILE, chroma_loader=build_chroma_loader(CHROMA_DIR))

def save_state(state: StateV2):
    state.save(STATE_FILE)

def fingerprint(path: str, mtime: float) -> str:
    """[v2] path+mtime 指纹（保留兼容，state v2 旧条目仍用此判定）。

    历史 bug：mtime 变化即使内容不变也触发重处理，且 chroma GC 后无法 recover。
    新代码请优先使用 content_fingerprint(path, content)。
    """
    return hashlib.sha256(f"{path}:{mtime:.3f}".encode()).hexdigest()[:32]


def content_fingerprint(path: str, content: bytes) -> str:
    """[v3] path+content 指纹（首选）。

    优势：mtime 变但内容不变 → fp 不变 → 跳过；
    chroma GC 后从 metadata 反向重建 dict[path → content_fp] 恒可行。
    """
    h = hashlib.sha256()
    h.update(path.encode("utf-8", errors="ignore"))
    h.update(b":")
    h.update(content)
    return h.hexdigest()[:32]

# ========== 文档解析 ==========
def read_file(path: Path) -> str:
    """读取文件内容，支持多格式"""
    ext = path.suffix.lower()
    try:
        if ext in ('.md', '.txt', '.sql', '.py', '.sh', '.yaml', '.yml'):
            return path.read_text(encoding='utf-8', errors='ignore')

        elif ext == '.pdf':
            from pdfminer.high_level import extract_text
            text = extract_text(str(path))
            return text or ''

        elif ext == '.docx':
            import docx
            doc = docx.Document(str(path))
            return '\n'.join(p.text for p in doc.paragraphs if p.text.strip())

        elif ext == '.xlsx':
            import openpyxl
            wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
            lines = []
            for ws in wb.worksheets:
                lines.append(f"Sheet: {ws.title}")
                for row in ws.iter_rows(max_row=200, values_only=True):
                    row_text = ' | '.join(str(c) for c in row if c is not None)
                    if row_text.strip():
                        lines.append(row_text)
            return '\n'.join(lines)

        elif ext == '.pptx':
            from pptx import Presentation
            prs = Presentation(str(path))
            lines = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        lines.append(shape.text)
            return '\n'.join(lines)

    except Exception as e:
        raise RuntimeError(f"解析失败({ext}): {e}")
    return ''

# ========== 扫描文件 ==========
def scan_files() -> list:
    """扫描所有支持的文件，返回 (path, mtime) 列表"""
    files = []
    for base in SCAN_DIRS:
        base_path = Path(base)
        if not base_path.exists():
            continue
        for f in base_path.rglob('*'):
            try:
                if f.suffix.lower() not in SUPPORTED_EXT:
                    continue
                if any(ex in f.parts for ex in EXCLUDE_DIRS):
                    continue
                if '.bak.' in f.name:
                    continue
                # 准入闸（loki_scan_filter）：拦第三方库/虚拟环境/缓存等脏数据
                if is_excluded_path(f):
                    continue
                # 文件名低价值过滤（防 33GB 内存爆事故 2026-04-28）：
                # 拦 test_*.py / *_test.py 等测试垃圾代码
                if is_low_value_file(f):
                    continue
                st = f.stat()
                if st.st_size < 50:
                    continue
                files.append((f, st.st_mtime))
            except (FileNotFoundError, PermissionError, OSError):
                pass  # 跳过失效symlink/无权限
    return files

# ========== Part 1: 文档向量化 ==========
def run_docs(model, col, state: StateV2, max_files: int | None = None):
    """返回 (state, model) — 并行读 + 批量化 embed 流水线。"""
    import concurrent.futures

    log.info("=" * 60)
    log.info("Loki Part 1: 文档向量化 (MD/PDF/Office/SQL)")

    all_files = scan_files()
    log.info(f"  扫描到文件: {len(all_files)}")

    todo = [(f, mt) for f, mt in all_files
            if not state.is_doc_done(str(f), fingerprint(str(f), mt))]
    skipped = len(all_files) - len(todo)
    log.info(f"  待处理: {len(todo)} | 跳过已有: {skipped}")

    if max_files is not None and len(todo) > max_files:
        log.info(f"  ⚠️  --max-files={max_files} 触发：本次只处理前 {max_files} 个，"
                 f"剩余 {len(todo) - max_files} 个待下次跑")
        todo = todo[:max_files]

    # ---- Phase 1: 并行读文件 + content-skip 判定 ----
    # 用线程池把 I/O-bound 的文件读取从 GPU embed 循环中剥离，
    # 8 个 worker 并行读盘，M5 NVMe 基本喂得饱 GPU
    success = failed = content_skip = 0
    ready: list[tuple[Path, float, str, str]] = []  # (f, mtime, content_fp, body)

    def _read_one(item: tuple[Path, float]) -> tuple:
        f, mt = item
        try:
            text = read_file(f).strip()
            fp = content_fingerprint(str(f), text.encode("utf-8", errors="ignore"))
            return (f, mt, text, fp, None)
        except Exception as e:
            return (f, mt, None, None, str(e))

    with concurrent.futures.ThreadPoolExecutor(max_workers=8) as pool:
        for f, mt, text, content_fp, err in pool.map(_read_one, todo):
            if err:
                log.error(f"  FAIL {f.name}: {err}")
                failed += 1
                state.mark_doc(str(f), fingerprint(str(f), mt))
                continue

            mtime_fp = fingerprint(str(f), mt)
            if state.is_doc_done(str(f), content_fp, mtime_fp):
                state.mark_doc(str(f), content_fp)
                content_skip += 1
                continue

            if len(text) < 30:
                log.info(f"  SKIP(空) {f.name}")
                state.mark_doc(str(f), content_fp)
                content_skip += 1
                continue

            body = extract_for_embed(f, text)
            ready.append((f, mt, content_fp, body))

    # ---- Phase 2: 批量 embed + chroma 入库 ----
    for i in range(0, len(ready), BATCH_SIZE):
        batch = ready[i:i+BATCH_SIZE]
        ids, docs, metas, paths = [], [], [], []

        for f, mt, content_fp, body in batch:
            embed_text = f"文件:{f.name}\n{body}"
            ids.append(content_fp)
            docs.append(embed_text)
            metas.append({
                'path': str(f),
                'name': f.name,
                'ext': f.suffix.lower(),
                'mtime': mt,
                'content_fp': content_fp,
                'source': 'loki_docs',
            })
            paths.append(str(f))
            log.info(f"  [{i+len(ids)}/{len(ready)}] {f.suffix} {f.name}")

        if not ids:
            continue

        try:
            seen, u_ids, u_docs, u_metas, u_paths = set(), [], [], [], []
            for id_, doc_, meta_, p_ in zip(ids, docs, metas, paths):
                if id_ not in seen:
                    seen.add(id_)
                    u_ids.append(id_); u_docs.append(doc_); u_metas.append(meta_); u_paths.append(p_)
            ids, docs, metas, paths = u_ids, u_docs, u_metas, u_paths

            embeddings = model.encode(docs, normalize_embeddings=True,
                                      show_progress_bar=False).tolist()
            col.upsert(ids=ids, documents=docs, metadatas=metas, embeddings=embeddings)
            for path_, fp_item in zip(paths, ids):
                state.mark_doc(path_, fp_item)
            success += len(ids)
            save_state(state)
            log.info(f"  ✅ 批次入库 {len(ids)} 条 (总计 {col.count()})")
        except Exception as e:
            log.error(f"  FAIL 入库: {e}")
            failed += len(ids)
        finally:
            try:
                del embeddings
            except NameError:
                pass
            del ids, docs, metas, paths
            _free_mps_memory()

        if SLEEP_BATCH > 0:
            time.sleep(SLEEP_BATCH)

    log.info(f"Part 1 完成: 成功={success} 失败={failed} content跳过={content_skip} 库总量={col.count()}")
    return state, model

# ========== Part 2: DDL向量化 ==========
def run_ddl(model, col, state: StateV2) -> StateV2:
    log.info("=" * 60)
    log.info("Loki Part 2: MySQL DDL 向量化")

    try:
        conn = pymysql.connect(**DB_CONFIG, database='information_schema')
        cur = conn.cursor()
        cur.execute("""
            SELECT TABLE_SCHEMA, TABLE_NAME, TABLE_COMMENT
            FROM TABLES
            WHERE TABLE_SCHEMA NOT IN
              ('information_schema','performance_schema','mysql','sys')
            ORDER BY TABLE_SCHEMA, TABLE_NAME
        """)
        tables = cur.fetchall()
        log.info(f"  发现 {len(tables)} 张表")
    except Exception as e:
        log.error(f"  MySQL连接失败: {e}")
        return state

    # 过滤备份/归档/临时表（详见 loki_ddl_filter）
    before_filter = len(tables)
    tables = [(db, tbl, tc) for db, tbl, tc in tables if not is_excluded_table(db, tbl)]
    if before_filter != len(tables):
        log.info(f"  过滤备份/归档/临时表: -{before_filter - len(tables)} 张")

    todo = []
    for db, tbl, tcomment in tables:
        ddl_id = hashlib.md5(f"ddl:{db}.{tbl}".encode()).hexdigest()
        if not state.is_ddl_done(ddl_id):
            todo.append((db, tbl, tcomment, ddl_id))

    log.info(f"  待处理: {len(todo)} | 跳过已有: {len(tables)-len(todo)}")

    success = failed = 0
    for i in range(0, len(todo), BATCH_SIZE):
        batch = todo[i:i+BATCH_SIZE]
        ids, docs, metas = [], [], []

        for db, tbl, tcomment, ddl_id in batch:
            try:
                cur.execute(f"""
                    SELECT COLUMN_NAME, COLUMN_TYPE, IS_NULLABLE,
                           COLUMN_DEFAULT, COLUMN_COMMENT
                    FROM COLUMNS
                    WHERE TABLE_SCHEMA='{db}' AND TABLE_NAME='{tbl}'
                    ORDER BY ORDINAL_POSITION
                """)
                cols = cur.fetchall()
                lines = [f"数据库:{db}", f"表:{tbl}"]
                if tcomment:
                    lines.append(f"说明:{tcomment}")
                lines.append("字段:")
                for cname, ctype, nullable, default, ccomment in cols:
                    line = f"  {cname} {ctype}"
                    if ccomment:
                        line += f" -- {ccomment}"
                    lines.append(line)

                ids.append(ddl_id)
                docs.append('\n'.join(lines)[:MAX_CHARS])
                metas.append({'db': db, 'table': tbl, 'comment': tcomment or '',
                               'col_count': len(cols), 'source': 'ddl'})
                log.info(f"  [{i+len(ids)}/{len(todo)}] {db}.{tbl} ({len(cols)}字段)")
            except Exception as e:
                log.error(f"  FAIL {db}.{tbl}: {e}")
                failed += 1

        if not ids:
            continue
        try:
            embeddings = model.encode(docs, normalize_embeddings=True,
                                      show_progress_bar=False).tolist()
            col.upsert(ids=ids, documents=docs, metadatas=metas, embeddings=embeddings)
            for ddl_id in ids:
                state.mark_ddl(ddl_id)
            success += len(ids)
            save_state(state)
            log.info(f"  ✅ 批次入库 {len(ids)} 条 (DDL总量={col.count()})")
        except Exception as e:
            log.error(f"  FAIL 入库: {e}")
            failed += len(ids)
        finally:
            try:
                del embeddings
            except NameError:
                pass
            del ids, docs, metas
            _free_mps_memory()

        if SLEEP_BATCH > 0:
            time.sleep(SLEEP_BATCH)

    conn.close()
    log.info(f"Part 2 完成: 成功={success} 失败={failed} DDL总量={col.count()}")
    return state

# ========== Part 2: 长文档分块 ==========

CHUNK_EXTS = {'.md', '.pdf', '.docx', '.txt'}
CHUNK_MIN_SIZE = 3000
CHUNK_SKIP_SIZE = 5_000_000

def _is_chunk_candidate(f: Path) -> bool:
    if f.suffix.lower() not in CHUNK_EXTS:
        return False
    try:
        st = f.stat()
        if st.st_size < CHUNK_MIN_SIZE:
            return False
        if st.st_size > CHUNK_SKIP_SIZE:
            return False
        name = f.name
        if name == '外部笔记索引.md':
            return False
        if name.startswith('cursor_') and st.st_size > 10_000_000:
            return False
        if name.startswith('backup_') and name.endswith('.sql'):
            return False
    except OSError:
        return False
    return True


def run_chunks(model, col, state: StateV2):
    import concurrent.futures
    from loki_chunk import get_chunks, chunk_pdf_pages

    log.info("=" * 60)
    log.info("Part 2: 长文档分块向量化 (>3KB 文件按语义切块)")

    all_files = scan_files()
    candidates = [(f, mt) for f, mt in all_files if _is_chunk_candidate(f)]
    log.info(f"  扫描到可切块文件: {len(candidates)}")

    todo = []
    success = failed = content_skip = 0

    def _read_one(item):
        f, mt = item
        try:
            text = read_file(f).strip() if f.suffix.lower() != '.pdf' else ''
            fp = content_fingerprint(str(f), text.encode("utf-8", errors="ignore"))
            return (f, mt, text, fp, None)
        except Exception as e:
            return (f, mt, None, None, str(e))

    with concurrent.futures.ThreadPoolExecutor(max_workers=8) as pool:
        for f, mt, text, content_fp, err in pool.map(_read_one, candidates):
            if err:
                log.error(f"  FAIL {f.name}: {err}")
                failed += 1
                continue
            if state.is_chunk_done(str(f), content_fp):
                content_skip += 1
                continue
            todo.append((f, content_fp, text))

    log.info(f"  待切块: {len(todo)} | 跳过已有: {content_skip} | 失败: {failed}")

    all_ids, all_docs, all_metas = [], [], []
    total_chunks = 0

    for file_idx, (f, fp, text) in enumerate(todo):
        try:
            if f.suffix.lower() == '.pdf':
                raw = read_file(f).strip()
                chunks_data = chunk_pdf_pages(f) if raw else []
            else:
                raw = text or read_file(f).strip()
                chunks_data = get_chunks(f, raw)

            if not chunks_data:
                state.mark_chunk(str(f), fp)
                continue

            rel_path = str(f)
            for base in SCAN_DIRS:
                if rel_path.startswith(base):
                    rel_path = rel_path[len(base):].lstrip('/')
                    break

            for ci, chunk in enumerate(chunks_data):
                chunk_id = f"{fp}_chunk_{ci:03d}"
                chunk_text = chunk['text']
                embed_text = f"文件:{f.name}\n块:{ci+1}/{len(chunks_data)}\n---\n{chunk_text[:1500]}"
                all_ids.append(chunk_id)
                all_docs.append(embed_text)
                all_metas.append({
                    'source_file': rel_path,
                    'chunk_index': ci,
                    'total_chunks': len(chunks_data),
                    'heading': chunk.get('heading', '')[:200],
                    'ext': f.suffix.lower(),
                    'file_id': fp,
                    'type': chunk.get('type', 'narrative'),
                    'char_count': len(chunk_text),
                    'source': 'loki_chunk',
                })

            state.mark_chunk(str(f), fp)
            success += 1
            total_chunks += len(chunks_data)
            log.info(f"  [{file_idx+1}/{len(todo)}] {f.name} -> {len(chunks_data)} chunks")
        except Exception as e:
            log.error(f"  FAIL {f.name}: {e}")
            failed += 1

    if all_ids:
        for bi in range(0, len(all_ids), BATCH_SIZE):
            b_ids = all_ids[bi:bi+BATCH_SIZE]
            b_docs = all_docs[bi:bi+BATCH_SIZE]
            b_metas = all_metas[bi:bi+BATCH_SIZE]
            try:
                embeddings = model.encode(b_docs, normalize_embeddings=True,
                                          show_progress_bar=False).tolist()
                col.upsert(ids=b_ids, documents=b_docs, metadatas=b_metas,
                          embeddings=embeddings)
                log.info(f"  ✅ chunk 批次入库 {len(b_ids)} 条 (chunk 库={col.count()})")
            except Exception as e:
                log.error(f"  FAIL chunk 入库 batch {bi}: {e}")
                failed += len(b_ids)
            finally:
                _free_mps_memory()

    _gc_chunk_orphans(col, state)
    log.info(f"Part 2 完成: 文件={success} 失败={failed} 总chunk={total_chunks} 库总量={col.count()}")
    return state


def _gc_chunk_orphans(col, state):
    total = col.count()
    if total == 0:
        return

    valid_file_ids = set(fp for fp in state.chunks.values())

    all_ids = []
    offset = 0
    while offset < total:
        result = col.get(limit=1000, offset=offset, include=[])
        ids = result.get('ids') or []
        if not ids:
            break
        all_ids.extend(ids)
        offset += 1000

    orphans = [cid for cid in all_ids
               if '_chunk_' in cid and cid.split('_chunk_')[0] not in valid_file_ids]

    if orphans:
        for i in range(0, len(orphans), 100):
            col.delete(ids=orphans[i:i+100])
        log.info(f"  GC: 已删除 {len(orphans)} 条孤儿 chunk (库剩余 {col.count()})")
    else:
        log.info(f"  GC: 无孤儿 chunk")


# ========== BM25 索引 rebuild ==========
BM25_INDEX_PATH = LOG_DIR / "bm25_index.pkl"
USERDICT_PATH   = Path(__file__).parent / "loki_userdict.txt"

def rebuild_bm25(client):
    """从 ChromaDB 全量 rebuild BM25 索引（增量 tokenize 缓存）

    优化：上一轮 pickle 中按 (doc_id, sig) 缓存 tokens，
    本轮命中则直接复用，未命中才走 jieba.cut。稳态命中率 >95%，
    188min 预处理 → 几十秒。详见 bm25_tokens_cache.py / OB 4/25 诊断笔记。
    """
    import pickle, jieba
    from rank_bm25 import BM25Okapi
    from bm25_tokens_cache import compute_sig, load_cache, lookup

    t_start = time.time()
    log.info("=" * 60)
    log.info("Loki Phase 3: BM25 索引 rebuild")

    jieba.load_userdict(str(USERDICT_PATH))
    stopwords = set("的 了 是 在 和 有 就 不 也 人 都 一 个 上 我 中 到 大 为 这 与 他 它 要 会 可以 没有 对 对于".split())

    def tokenize(text):
        return [w.lower().strip() for w in jieba.cut(text) if w.strip() and w.strip() not in stopwords]

    cache = load_cache(BM25_INDEX_PATH)
    log.info(f"  上一轮 tokens 缓存: {len(cache)} 条")

    collections_map = {
        'summaries': 'doc_knowledge_bge_m3',
        'chunks':    'doc_knowledge_chunks',
        'ddl':       'ddl_schema_bge_m3',
    }

    corpus_tokens = []
    doc_ids = []
    doc_meta = []
    doc_text_sigs = []
    total_hits = 0
    total_misses = 0

    for source_key, col_name in collections_map.items():
        try:
            col = client.get_collection(col_name)
            count = col.count()
            if count == 0:
                continue
        except Exception:
            continue

        t_col = time.time()
        offset = 0
        batch_size = 5000
        loaded = 0
        col_hits = 0
        while offset < count:
            results = col.get(limit=batch_size, offset=offset, include=['documents', 'metadatas'])
            if not results['ids']:
                break
            for i, doc_id in enumerate(results['ids']):
                doc_text = results['documents'][i] or ''
                meta = results['metadatas'][i] or {}
                searchable = doc_text
                for field in ['source_file', 'name', 'path', 'heading', 'topic', 'domain', 'table', 'db', 'description']:
                    if field in meta and meta[field]:
                        searchable += ' ' + str(meta[field])
                sig = compute_sig(searchable)
                cached = lookup(cache, doc_id, sig)
                if cached is not None:
                    tokens = cached
                    col_hits += 1
                else:
                    tokens = tokenize(searchable)
                if not tokens:
                    continue
                corpus_tokens.append(tokens)
                doc_ids.append(doc_id)
                doc_text_sigs.append(sig)
                file_key = meta.get('source_file', meta.get('path', meta.get('name', meta.get('table', doc_id))))
                doc_meta.append({
                    'source': source_key, 'file': file_key,
                    'heading': meta.get('heading', ''), 'topic': meta.get('topic', meta.get('table', '')),
                    'domain': meta.get('domain', source_key), 'doc_preview': doc_text[:500],
                })
                loaded += 1
            offset += batch_size
        total_hits += col_hits
        total_misses += (loaded - col_hits)
        hit_rate = (col_hits * 100 // loaded) if loaded else 0
        log.info(f"  {col_name}: {loaded} 条 (cache hit {col_hits}/{loaded} = {hit_rate}%, "
                 f"load+tokenize {time.time()-t_col:.1f}s)")

    if not corpus_tokens:
        log.warning("  无文档，跳过 BM25 rebuild")
        return

    log.info(f"  cache 总命中: {total_hits}/{total_hits+total_misses} "
             f"({total_hits*100//(total_hits+total_misses) if (total_hits+total_misses) else 0}%)")

    t_build = time.time()
    bm25 = BM25Okapi(corpus_tokens)
    log.info(f"  BM25Okapi 构造: {time.time()-t_build:.1f}s")

    t_pickle = time.time()
    index_data = {
        'bm25': bm25, 'doc_ids': doc_ids, 'doc_meta': doc_meta,
        'corpus_tokens': corpus_tokens,
        'doc_text_sigs': doc_text_sigs,  # H.2: 增量缓存 key
        'build_time': datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
        'doc_count': len(doc_ids),
    }
    with open(BM25_INDEX_PATH, 'wb') as f:
        pickle.dump(index_data, f)

    size_mb = BM25_INDEX_PATH.stat().st_size / 1024 / 1024
    log.info(f"  pickle 写盘: {time.time()-t_pickle:.1f}s ({size_mb:.1f}MB)")
    log.info(f"  BM25 索引 rebuild 完成: {len(doc_ids)} 条, 总耗时 {time.time()-t_start:.1f}s")


# ========== 主流程 ==========
def _parse_args(argv: list) -> tuple[str, int | None]:
    """解析 mode + --max-files。保持向后兼容（位置参数 mode）。"""
    mode = 'all'
    max_files: int | None = None
    i = 1
    while i < len(argv):
        a = argv[i]
        if a.startswith("--max-files="):
            max_files = int(a.split("=", 1)[1])
        elif a == "--max-files" and i + 1 < len(argv):
            max_files = int(argv[i + 1])
            i += 1
        elif not a.startswith("--"):
            mode = a
        i += 1
    return mode, max_files


def main():
    mode, max_files = _parse_args(sys.argv)
    log.info("=" * 60)
    log.info(f"Loki 启动 mode={mode} max_files={max_files}  "
             f"{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")

    # 1. 先加载模型（无副作用，不需要锁）
    log.info("加载 BGE-M3...")
    model = SentenceTransformer(BGE_PATH)
    log.info(f"  维度: {model.get_sentence_embedding_dimension()}")

    # 2. 获取互斥锁（拿不到就干净退出，不会泄漏资源）
    LOCK_FILE.parent.mkdir(parents=True, exist_ok=True)
    lock_fp = open(LOCK_FILE, 'w')
    try:
        fcntl.flock(lock_fp, fcntl.LOCK_EX | fcntl.LOCK_NB)
        lock_fp.write(f"pipeline:{os.getpid()}\n"); lock_fp.flush()
        log.info("🔒 已获取ChromaDB写锁")
    except (IOError, OSError):
        log.error("❌ 无法获取写锁，reindex/chunk可能正在运行。退出。")
        lock_fp.close()
        sys.exit(1)

    # 3. 在锁保护下初始化ChromaDB + 做工作，finally保证释放
    try:
        client = chromadb.PersistentClient(CHROMA_DIR)
        col_docs = client.get_or_create_collection(COLLECTION_DOCS,
                       metadata={"hnsw:space": "cosine"})
        col_ddl  = client.get_or_create_collection(COLLECTION_DDL,
                       metadata={"hnsw:space": "cosine"})
        col_chunks = client.get_or_create_collection(COLLECTION_CHUNKS,
                         metadata={"hnsw:space": "cosine"})

        state = load_state()
        log.info(f"  已处理状态: docs={len(state.docs)} chunks={len(state.chunks)} ddl={len(state.ddl)}")

        t0 = time.time()
        if mode in ('all', 'docs'):
            state, model = run_docs(model, col_docs, state, max_files=max_files)
        if mode in ('all', 'chunks'):
            state = run_chunks(model, col_chunks, state)
        if mode in ('all', 'ddl'):
            state = run_ddl(model, col_ddl, state)

        save_state(state)

        # Phase 3: BM25 索引 rebuild（每轮都更新，<10s）
        rebuild_bm25(client)

        elapsed = (time.time() - t0) / 60
        log.info("=" * 60)
        log.info(f"Loki 完成 耗时={elapsed:.1f}分钟")
        log.info(f"  文档库: {col_docs.count()} 条")
        log.info(f"  Chunk库: {col_chunks.count()} 条")
        log.info(f"  DDL库:  {col_ddl.count()} 条")
        log.info(f"  日志:   {log_file}")
    finally:
        fcntl.flock(lock_fp, fcntl.LOCK_UN); lock_fp.close()
        log.info("🔓 已释放ChromaDB写锁")

if __name__ == '__main__':
    main()
