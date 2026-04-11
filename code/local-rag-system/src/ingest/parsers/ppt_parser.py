from typing import List, Dict, Any, Optional
from pathlib import Path
import re
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .base_parser import BaseParser
from ...models.document import DocumentWithChunks, DocumentChunkWithMetadata, PPTMetadata, DataSource

class PPTParser(BaseParser):
    """PowerPoint文件解析器"""
    
    def __init__(self, config: Dict[str, Any] = None):
        super().__init__(config)
        self.extract_notes = self.config.get('extract_notes', True)
        self.combine_title_body = self.config.get('combine_title_body', True)
        self.max_tokens_per_slide = self.config.get('max_tokens_per_slide', 800)
    
    def can_parse(self, file_path: str) -> bool:
        """检查是否为PowerPoint文件"""
        return Path(file_path).suffix.lower() in ['.pptx']
    
    def parse(self, file_path: str) -> DocumentWithChunks:
        """解析PowerPoint文件"""
        doc_id = self._generate_doc_id(file_path)
        file_size = self._get_file_size(file_path)
        file_hash = self._calculate_file_hash(file_path)
        
        # 创建文档对象
        document = DocumentWithChunks(
            doc_id=doc_id,
            file_path=file_path,
            file_type=DataSource.PPT,
            file_size_bytes=file_size,
            sha256_hash=file_hash
        )
        
        try:
            # 读取PowerPoint文件
            chunks = self._parse_presentation(file_path, doc_id)
            document.chunks = chunks
            document.processing_status = "completed"
            
        except Exception as e:
            document.processing_status = "failed"
            print(f"Error parsing PPT file {file_path}: {e}")
            raise
        
        return document
    
    def _parse_presentation(self, file_path: str, doc_id: str) -> List[DocumentChunkWithMetadata]:
        """解析PowerPoint演示文稿"""
        chunks = []
        
        try:
            prs = Presentation(file_path)
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_chunks = self._parse_slide(slide, slide_num, doc_id)
                chunks.extend(slide_chunks)
                
        except Exception as e:
            print(f"Failed to parse presentation: {e}")
            raise
        
        return chunks
    
    def _parse_slide(self, slide, slide_num: int, doc_id: str) -> List[DocumentChunkWithMetadata]:
        """解析单个幻灯片"""
        slide_content = self._extract_slide_content(slide)
        
        if not slide_content.strip():
            return []
        
        # 创建PPT元数据
        slide_title = self._extract_slide_title(slide)
        has_notes = self._extract_notes_content(slide) != ""
        
        ppt_metadata = PPTMetadata(
            slide_number=slide_num,
            slide_title=slide_title,
            has_notes=has_notes
        )
        
        # 检查是否需要二次切分
        estimated_tokens = len(slide_content.split())
        
        if estimated_tokens <= self.max_tokens_per_slide:
            # 不需要切分，直接创建一个块
            chunk = DocumentChunkWithMetadata(
                doc_id=doc_id,
                content=slide_content,
                chunk_index=0,
                token_count=estimated_tokens,
                metadata={
                    "slide_number": slide_num,
                    "slide_title": slide_title,
                    "has_notes": has_notes
                },
                ppt_metadata=ppt_metadata
            )
            return [chunk]
        else:
            # 需要二次切分
            return self._split_slide_content(
                slide_content, slide_num, doc_id, ppt_metadata
            )
    
    def _extract_slide_content(self, slide) -> str:
        """提取幻灯片内容"""
        content_parts = []
        
        # 提取标题
        if slide.shapes.title:
            title_text = slide.shapes.title.text.strip()
            if title_text:
                content_parts.append(f"标题: {title_text}")
        
        # 提取正文内容
        body_content = []
        
        for shape in slide.shapes:
            # 跳过标题形状（已处理）
            if shape == slide.shapes.title:
                continue
                
            # 处理文本框
            if hasattr(shape, "text_frame") and shape.text_frame:
                shape_text = self._extract_text_from_shape(shape)
                if shape_text.strip():
                    body_content.append(shape_text)
            
            # 处理表格
            elif hasattr(shape, "table"):
                table_text = self._extract_text_from_table(shape)
                if table_text.strip():
                    body_content.append(table_text)
        
        # 添加正文内容
        if body_content:
            if self.combine_title_body and content_parts:
                # 标题和正文合并
                content_parts.extend(body_content)
            else:
                # 正文独立段落
                for body_part in body_content:
                    content_parts.append(body_part)
        
        # 提取备注
        if self.extract_notes:
            notes_content = self._extract_notes_content(slide)
            if notes_content.strip():
                content_parts.append(f"备注: {notes_content}")
        
        # 连接所有内容
        if content_parts:
            if self.combine_title_body:
                return "\n".join(content_parts)
            else:
                return "\n---\n".join(content_parts)
        else:
            return ""
    
    def _extract_text_from_shape(self, shape) -> str:
        """从形状中提取文本"""
        try:
            if not shape.text_frame:
                return ""
            
            text_parts = []
            for paragraph in shape.text_frame.paragraphs:
                paragraph_text = paragraph.text.strip()
                if paragraph_text:
                    text_parts.append(paragraph_text)
            
            return "\n".join(text_parts)
        except Exception:
            return ""
    
    def _extract_text_from_table(self, table) -> str:
        """从表格中提取文本"""
        try:
            table_rows = []
            
            for row in table.rows:
                row_cells = []
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    if cell_text:
                        row_cells.append(cell_text)
                    else:
                        row_cells.append("")
                
                if row_cells:
                    table_rows.append(" | ".join(row_cells))
            
            if table_rows:
                return "表格内容:\n" + "\n".join(table_rows)
            else:
                return ""
        except Exception:
            return ""
    
    def _extract_slide_title(self, slide) -> Optional[str]:
        """提取幻灯片标题"""
        try:
            if slide.shapes.title:
                return slide.shapes.title.text.strip()
            else:
                return None
        except Exception:
            return None
    
    def _extract_notes_content(self, slide) -> str:
        """提取幻灯片备注内容"""
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                return slide.notes_slide.notes_text_frame.text.strip()
            else:
                return ""
        except Exception:
            return ""
    
    def _split_slide_content(self, content: str, slide_num: int, 
                            doc_id: str, ppt_metadata: PPTMetadata) -> List[DocumentChunkWithMetadata]:
        """二次切分幻灯片内容"""
        # 按段落分割内容
        paragraphs = content.split('\n')
        
        chunks = []
        current_chunk_content = []
        current_token_count = 0
        chunk_index = 0
        
        for paragraph in paragraphs:
            paragraph_tokens = len(paragraph.split())
            
            # 如果单个段落过长，需要进一步分割
            if paragraph_tokens > self.max_tokens_per_slide:
                # 先保存当前块（如果有内容）
                if current_chunk_content:
                    chunk = self._create_ppt_chunk(
                        current_chunk_content, slide_num, doc_id, 
                        ppt_metadata, chunk_index, current_token_count
                    )
                    chunks.append(chunk)
                    chunk_index += 1
                    current_chunk_content = []
                    current_token_count = 0
                
                # 分割长段落
                sub_chunks = self._split_long_paragraph(paragraph, self.max_tokens_per_slide)
                for i, sub_chunk in enumerate(sub_chunks):
                    chunk = self._create_ppt_chunk(
                        [sub_chunk], slide_num, doc_id,
                        ppt_metadata, chunk_index, len(sub_chunk.split())
                    )
                    chunks.append(chunk)
                    chunk_index += 1
            else:
                # 检查是否超过块大小
                if current_token_count + paragraph_tokens > self.max_tokens_per_slide:
                    # 保存当前块
                    if current_chunk_content:
                        chunk = self._create_ppt_chunk(
                            current_chunk_content, slide_num, doc_id,
                            ppt_metadata, chunk_index, current_token_count
                        )
                        chunks.append(chunk)
                        chunk_index += 1
                        current_chunk_content = []
                        current_token_count = 0
                
                # 添加到当前块
                current_chunk_content.append(paragraph)
                current_token_count += paragraph_tokens
        
        # 处理最后剩余内容
        if current_chunk_content:
            chunk = self._create_ppt_chunk(
                current_chunk_content, slide_num, doc_id,
                ppt_metadata, chunk_index, current_token_count
            )
            chunks.append(chunk)
        
        return chunks
    
    def _split_long_paragraph(self, paragraph: str, max_tokens: int) -> List[str]:
        """分割长段落"""
        # 按句子分割
        sentences = re.split(r'[。！？；.!?;]', paragraph)
        
        chunks = []
        current_chunk = ""
        current_tokens = 0
        
        for sentence in sentences:
            sentence = sentence.strip()
            if not sentence:
                continue
            
            sentence_tokens = len(sentence.split())
            
            if current_tokens + sentence_tokens > max_tokens:
                if current_chunk:
                    chunks.append(current_chunk)
                    current_chunk = ""
                    current_tokens = 0
                
                # 如果单个句子过长，按字符强制分割
                if sentence_tokens > max_tokens:
                    char_chunks = self._split_by_characters(sentence, max_tokens)
                    chunks.extend(char_chunks)
                else:
                    current_chunk = sentence
                    current_tokens = sentence_tokens
            else:
                if current_chunk:
                    current_chunk += "。" + sentence
                else:
                    current_chunk = sentence
                current_tokens += sentence_tokens
        
        if current_chunk:
            chunks.append(current_chunk)
        
        return chunks
    
    def _split_by_characters(self, text: str, max_tokens: int) -> List[str]:
        """按字符强制分割文本"""
        # 估算每个汉字约等于1个token
        max_chars = max_tokens * 1.5  # 留一些buffer
        
        chunks = []
        for i in range(0, len(text), int(max_chars)):
            chunk = text[i:i + int(max_chars)]
            if chunk.strip():
                chunks.append(chunk)
        
        return chunks
    
    def _create_ppt_chunk(self, content_parts: List[str], slide_num: int, 
                         doc_id: str, ppt_metadata: PPTMetadata, 
                         chunk_index: int, token_count: int) -> DocumentChunkWithMetadata:
        """创建PPT文档块"""
        content = "\n".join(content_parts)
        
        chunk = DocumentChunkWithMetadata(
            doc_id=doc_id,
            content=content,
            chunk_index=chunk_index,
            token_count=token_count,
            metadata={
                "slide_number": slide_num,
                "slide_title": ppt_metadata.slide_title,
                "has_notes": ppt_metadata.has_notes,
                "chunk_part": f"{chunk_index + 1}"
            },
            ppt_metadata=ppt_metadata
        )
        
        return chunk
    
    def _generate_doc_id(self, file_path: str) -> str:
        """生成文档ID"""
        return f"ppt_{Path(file_path).stem}"